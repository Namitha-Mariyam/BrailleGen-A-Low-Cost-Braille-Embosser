import os, tempfile, cv2, numpy as np, easyocr, serial, time
from pptx import Presentation
from pdf2image import convert_from_path
from PIL import Image
from google.genai import Client
from google.genai.types import GenerateContentConfig
import re

# ---------------- CONFIG ----------------
GEMINI_API_KEY = "Your API key"
client = Client(api_key=GEMINI_API_KEY)
reader = easyocr.Reader(['en'], gpu=False)

# ---------------- FUNCTIONS ----------------

def gemini_extract(image_path):
    try:
        img = Image.open(image_path)
        response = client.models.generate_content(
            model="gemini-3.1-flash-lite-preview", 
            contents=["Extract the text and equations. Return equations in simple LaTeX format (e.g., \\frac{a}{b} and v^2). Do not use bold or italic markers.", img],
            config=GenerateContentConfig(temperature=0.1)
        )
        return response.text.strip()
    except Exception as e:
        print(f"[!] API Error: {e}. Using EasyOCR fallback.")
        results = reader.readtext(image_path, detail=0)
        return " ".join(results)
    
# ---------------- MAIN LOOP (CLEANED) ----------------

def main():
    path = input("Enter file path (Image/PDF/PPT): ").strip('"')
    if not os.path.exists(path): return

    ext = os.path.splitext(path.lower())[1]
    extracted = ""

    # 1. MULTI-FORMAT EXTRACTION
    if ext in [".jpg", ".jpeg", ".png"]:
        extracted = gemini_extract(path)
    elif ext == ".pdf":
        print("[INFO] Processing PDF...")
        pages = convert_from_path(path)
        for i, page in enumerate(pages):
            tmp = f"p_{i}.jpg"; page.save(tmp, "JPEG")
            extracted += gemini_extract(tmp) + "\n"
            os.remove(tmp)
    elif ext == ".pptx":
        print("[INFO] Processing PPT...")
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"): extracted += shape.text + "\n"

    print(f"\n[STEP 1] TEXT EXTRACTED:\n{extracted}")
    # ... (After extraction is complete) ...

    arduino = None
    try:
        # Update COM7 to your actual port if it changed
        arduino = serial.Serial('COM7', 9600, timeout=5)
        time.sleep(2) # Wait for Arduino to reboot after connection
        print("[INFO] Hardware connected.")
    except Exception as e:
        print(f"[SERIAL ERROR] Could not connect: {e}")
        return

    if arduino and arduino.is_open:
        # This is where the indentation fix happened
        lines = [l.strip() for l in extracted.split('\n') if l.strip()]
        
        for line in lines:
            print(f"[SENDING] {line}...")
            arduino.write((line + "\n").encode())
            
            # --- THE WAIT LOGIC (Handshake) ---
            print("[WAITING] Arduino is punching and moving Y-axis...")
            while True:
                if arduino.in_waiting > 0:
                    try:
                        response = arduino.readline().decode().strip()
                        print(f"[ARDUINO] {response}")
                        
                        if "Ready" in response:
                            print(f"[OK] Line complete. Sending next...")
                            break 
                    except:
                        pass # Ignore decoding glitches during motor movement
                
                time.sleep(0.1) 
            
        print("\n[SUCCESS] All lines printed. Closing connection.")
        arduino.close()

if __name__ == "__main__":
    main()
